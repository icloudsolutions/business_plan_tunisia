"""
PowerPoint presentation export for validated business plans.

generate_pptx_presentation(plan_data, audience, output_path)
audience: 'banque' | 'investisseur' | 'client'
"""

from __future__ import annotations

import shutil
import tempfile
from copy import copy
from datetime import datetime
from pathlib import Path
from typing import Any, Literal

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from bp_calc.engine import HORIZON, calculate_plan
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from tasks.export_excel import (
    ReportContext,
    _categorize_investments,
    _loan_annual_rows,
    _parse_plan_data,
    _product_names,
    _safe_div,
    build_plan_data,
)

Audience = Literal["banque", "investisseur", "client"]

# Design system
C_PRIMARY = RGBColor(0x00, 0x33, 0x66)
C_SECONDARY = RGBColor(0x33, 0x66, 0x99)
C_ACCENT = RGBColor(0xCC, 0x99, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0xF5, 0xF5, 0xF5)
C_CARD_BG = RGBColor(0xE8, 0xEE, 0xF5)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_SLIDE_NUM = RGBColor(0x99, 0x99, 0x99)

MPL_PRIMARY = "#003366"
MPL_SECONDARY = "#336699"
MPL_ACCENT = "#CC9900"
MPL_GREEN = "#336633"
MPL_ORANGE = "#CC6600"
MPL_RED = "#CC3333"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_H = Inches(0.07)

CHART_DPI = 300


def _mpl_readability() -> None:
    plt.rcParams.update(
        {
            "text.color": "#333333",
            "axes.labelcolor": "#333333",
            "axes.edgecolor": "#333333",
            "xtick.color": "#333333",
            "ytick.color": "#333333",
            "legend.labelcolor": "#333333",
            "figure.facecolor": "white",
            "axes.facecolor": "white",
        }
    )


def _style_pie_labels(ax) -> None:
    for t in ax.texts:
        t.set_color("#333333")
    for child in ax.get_children():
        if hasattr(child, "get_text") and child.get_text() and "%" in child.get_text():
            try:
                child.set_color("white")
                child.set_fontweight("bold")
            except Exception:
                pass


def _fmt(n: float | None, *, digits: int = 0) -> str:
    if n is None:
        return "—"
    return f"{n:,.{digits}f}".replace(",", " ")


def _pct(n: float | None) -> str:
    if n is None:
        return "—"
    return f"{n * 100:.1f}%"


def _drci_label(years: float | None) -> str:
    if years is None:
        return "—"
    y = int(years)
    m = int(round((years - y) * 12))
    if m:
        return f"{y} ans et {m} mois"
    return f"{y} ans"


def _audience_label(audience: Audience) -> str:
    return {"banque": "la Banque", "investisseur": "l'Investisseur", "client": "le Client"}[audience]


# —— Matplotlib charts (300 dpi PNG) ——


def _save_chart(fig, path: Path) -> str:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.tight_layout()
    fig.savefig(path, dpi=CHART_DPI, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return str(path)


def _chart_ca_bars(revenue: list[float], growth: list[float], path: Path) -> str:
    _mpl_readability()
    years = list(range(1, HORIZON + 1))
    fig, ax = plt.subplots(figsize=(10, 5))
    bars = ax.bar(years, revenue[:HORIZON], color=MPL_PRIMARY, alpha=0.9)
    ax.set_title("Chiffre d'affaires previsionnel (DT)", fontsize=12, color=MPL_PRIMARY, fontweight="bold")
    ax.set_xlabel("Annee")
    ax.set_ylabel("DT")
    ax.grid(True, axis="y", alpha=0.3)
    for bar, g in zip(bars, growth[:HORIZON]):
        h = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            h,
            f"{_fmt(h)}\n+{_pct(g)}" if g else _fmt(h),
            ha="center",
            va="bottom",
            fontsize=7,
            color="#333333",
        )
    return _save_chart(fig, path)


def _chart_results_lines(revenue, ebit, net, path: Path) -> str:
    _mpl_readability()
    years = list(range(1, HORIZON + 1))
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.plot(years, revenue, "o-", color=MPL_PRIMARY, label="CA HT", linewidth=2)
    ax.plot(years, ebit, "s-", color=MPL_SECONDARY, label="Resultat exploitation", linewidth=2)
    ax.plot(years, net, "^-", color=MPL_GREEN, label="Resultat net", linewidth=2)
    ax.axhline(0, color="#999", linewidth=0.8)
    ax.legend(fontsize=9)
    ax.grid(True, alpha=0.3)
    ax.set_title("Evolution des resultats", fontsize=12, color=MPL_PRIMARY, fontweight="bold")
    return _save_chart(fig, path)


def _chart_financing_pie(labels, values, path: Path) -> str:
    _mpl_readability()
    fig, ax = plt.subplots(figsize=(6, 6))
    colors = [MPL_PRIMARY, MPL_SECONDARY, MPL_ACCENT]
    ax.pie(
        values,
        labels=labels,
        autopct="%1.1f%%",
        colors=colors[: len(values)],
        startangle=90,
        textprops={"color": "#333333", "fontsize": 10},
    )
    _style_pie_labels(ax)
    ax.set_title("Schema de financement", fontsize=12, color=MPL_PRIMARY, fontweight="bold")
    return _save_chart(fig, path)


def _chart_cashflow_debt(ocf, loan_rows, path: Path) -> str:
    _mpl_readability()
    years = list(range(1, HORIZON + 1))
    closing = [r["closing"] for r in loan_rows]
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(years, ocf[:HORIZON], color=MPL_SECONDARY, alpha=0.85, label="CF exploitation")
    ax2 = ax.twinx()
    ax2.plot(years, closing, "o--", color=MPL_ORANGE, linewidth=2, label="Capital restant du")
    ax.set_title("Cash-flows et remboursement", fontsize=12, color=MPL_PRIMARY, fontweight="bold")
    ax.legend(loc="upper left", fontsize=8)
    ax2.legend(loc="upper right", fontsize=8)
    ax.grid(True, axis="y", alpha=0.3)
    return _save_chart(fig, path)


def _chart_cost_donut(values, labels, path: Path) -> str:
    _mpl_readability()
    fig, ax = plt.subplots(figsize=(7, 7))
    ax.pie(
        values,
        labels=labels,
        autopct="%1.1f%%",
        colors=[MPL_PRIMARY, MPL_TEAL, MPL_ORANGE, MPL_GREEN],
        pctdistance=0.75,
        startangle=90,
        textprops={"color": "#333333", "fontsize": 9},
    )
    _style_pie_labels(ax)
    centre = plt.Circle((0, 0), 0.55, fc="white")
    ax.add_artist(centre)
    ax.set_title("Structure des couts (moy. 7 ans)", fontsize=12, color=MPL_PRIMARY, fontweight="bold")
    return _save_chart(fig, path)


MPL_TEAL = "#336699"


def _chart_ipc_mini(ipc_data: dict, path: Path) -> str:
    _mpl_readability()
    years = ipc_data.get("years") or [2018, 2019, 2020, 2021, 2022]
    cats = ipc_data.get("categories") or {"Secteur": [3, 3.5, 4, 4.2, 4.5]}
    fig, ax = plt.subplots(figsize=(5, 3))
    for name, vals in list(cats.items())[:2]:
        ax.plot(years[: len(vals)], vals[: len(years)], marker="o", label=name)
    ax.legend(fontsize=7)
    ax.set_title("IPC secteur", fontsize=9, color=MPL_PRIMARY)
    ax.grid(True, alpha=0.3)
    return _save_chart(fig, path)


def _chart_investment_bars(items: list[tuple[str, float]], total: float, path: Path) -> str:
    _mpl_readability()
    names = [x[0][:25] for x in items]
    vals = [x[1] for x in items]
    fig, ax = plt.subplots(figsize=(9, max(3, len(items) * 0.4)))
    y_pos = range(len(names))
    ax.barh(y_pos, vals, color=MPL_SECONDARY)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(names, fontsize=8)
    for i, v in enumerate(vals):
        ax.text(v, i, f" {_fmt(v)} ({_pct(_safe_div(v, total))})", va="center", fontsize=8, color="#333333")
    ax.set_title("Plan d'investissement (DT)", fontsize=11, color=MPL_PRIMARY, fontweight="bold")
    return _save_chart(fig, path)


def _sensitivity_mini(ctx: ReportContext) -> list[list[str]]:
    """3x3 VAN table for ±10% CA and costs."""
    factors = [0.9, 1.0, 1.1]
    base = ctx.inputs.model_dump()
    rows = []
    for cost_f in factors:
        row = []
        for ca_f in factors:
            data = copy(base)
            ops = data.setdefault("operations", {})
            ops["salePrice"] = float(ops.get("salePrice", 0)) * ca_f
            ops["rawMaterialCost"] = float(ops.get("rawMaterialCost", 0)) * cost_f
            ops["packagingCost"] = float(ops.get("packagingCost", 0)) * cost_f
            from bp_schema.liasse import PlanInputs

            res = calculate_plan(PlanInputs.model_validate(data), discount_rate=ctx.ind.discountRate)
            row.append(_fmt(res.indicators.van))
        rows.append(row)
    return rows


class PptxDeck:
    def __init__(self, plan_data: dict, audience: Audience):
        self.plan_data = plan_data
        self.audience = audience
        self.ctx = _parse_plan_data(plan_data)
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._slide_no = 0
        self.market = plan_data.get("market_study") or {}
        self.swot = plan_data.get("swot") or {}
        self.promoter = plan_data.get("promoter") or {}
        self.team = plan_data.get("team") or []
        self.sector = str(plan_data.get("sector") or "Activite")
        self.site = str(plan_data.get("site") or "Tunisie")
        self.logo = plan_data.get("logo_path")
        self.chart_dir: Path | None = None
        self._figure_no = 0

    def _next_figure(self) -> int:
        self._figure_no += 1
        return self._figure_no

    def _figure_caption(self, slide, text: str, left, top, *, width=Inches(12)) -> None:
        n = self._next_figure()
        self._body(
            slide,
            f"Figure {n} — {text} (source : moteur bp_calc, donnees du plan)",
            left,
            top,
            width,
            Inches(0.45),
            size=10,
            italic=True,
            color=C_SLIDE_NUM,
        )

    def _table_caption(self, slide, text: str, left=Inches(0.4), top=None, *, width=Inches(12)) -> None:
        if top is None:
            top = Inches(6.5)
        self._body(
            slide,
            f"Tableau — {text} (source : moteur bp_calc)",
            left,
            top,
            width,
            Inches(0.35),
            size=10,
            italic=True,
            color=C_SLIDE_NUM,
        )

    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._slide_no += 1
        return slide

    def _footer(self, slide) -> None:
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H
        )
        band.fill.solid()
        band.fill.fore_color.rgb = C_PRIMARY
        band.line.fill.background()
        num = slide.shapes.add_textbox(SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.35), Inches(0.6), Inches(0.25))
        tf = num.text_frame
        tf.text = str(self._slide_no)
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = C_SLIDE_NUM
        p.alignment = PP_ALIGN.RIGHT
        if self.logo and Path(self.logo).is_file():
            try:
                slide.shapes.add_picture(self.logo, Inches(0.2), SLIDE_H - Inches(0.55), height=Inches(0.35))
            except Exception:
                pass

    def _title_box(self, slide, text: str, top=Inches(0.35), size=28) -> None:
        box = slide.shapes.add_textbox(Inches(0.5), top, SLIDE_W - Inches(1), Inches(0.7))
        tf = box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(size)
        p.font.color.rgb = C_PRIMARY
        p.font.name = "Calibri"

    def _body(self, slide, text: str, left, top, w, h, *, size=16, color=None, bold=False, italic=False):
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.bold = bold
        p.font.italic = italic
        if color:
            p.font.color.rgb = color
        else:
            p.font.color.rgb = C_TEXT
        return box

    def _add_picture(self, slide, path: str, left, top, width) -> None:
        if Path(path).is_file():
            slide.shapes.add_picture(path, left, top, width=width)

    def _fill_kpi_card(self, shape, icon: str, label: str, value: str) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_PRIMARY
        shape.line.width = Pt(1.25)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.clear()
        lines = [
            (icon, Pt(20), False, C_SECONDARY),
            (label, Pt(11), False, C_TEXT),
            (value, Pt(20), True, C_PRIMARY),
        ]
        for i, (text, size, bold, color) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.name = "Calibri"
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

    def _kpi_card(self, slide, left, top, w, h, icon: str, label: str, value: str) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, w, h)
        self._fill_kpi_card(shape, icon, label, value)

    def slide01_cover(self) -> None:
        slide = self._blank_slide()
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_PRIMARY
        bg.line.fill.background()
        if self.logo and Path(self.logo).is_file():
            try:
                slide.shapes.add_picture(
                    self.logo, SLIDE_W / 2 - Inches(1), Inches(1.2), width=Inches(2)
                )
            except Exception:
                pass
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), SLIDE_W - Inches(1), Inches(1))
        tf = title.text_frame
        tf.text = self.ctx.company
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), SLIDE_W - Inches(1), Inches(0.6))
        sub.text_frame.text = f"{self.sector} — {self.site}"
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = C_WHITE
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        foot = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), SLIDE_W - Inches(1), Inches(1))
        foot.text_frame.text = (
            f"{datetime.now().strftime('%d/%m/%Y')}\n"
            f"Promoteur : {self.promoter.get('name', self.ctx.company)}\n"
            f"Presente a {_audience_label(self.audience)}"
        )
        for p in foot.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER
        self._slide_no = 1

    def slide02_agenda(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Sommaire")
        chapters = [
            "1. Presentation du promoteur",
            "2. Le projet en un coup d'oeil",
            "3. Etude de marche",
            "4. Analyse SWOT",
            "5. Investissement et financement",
            "6. Performance financiere",
            "7. Conclusion",
        ]
        y = Inches(1.3)
        for ch in chapters:
            sq = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), y + Inches(0.08), Inches(0.15), Inches(0.15)
            )
            sq.fill.solid()
            sq.fill.fore_color.rgb = C_PRIMARY
            sq.line.fill.background()
            self._body(slide, ch, Inches(0.95), y, Inches(10), Inches(0.45), size=18)
            y += Inches(0.55)
        self._footer(slide)

    def slide03_promoter(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Presentation du promoteur")
        ph = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2.5), Inches(3)
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        self._body(slide, "[Photo]", Inches(1.2), Inches(2.5), Inches(1), Inches(0.4), size=14, color=C_SLIDE_NUM)
        info = (
            f"Nom : {self.promoter.get('name', self.ctx.company)}\n\n"
            f"Experience :\n{self.promoter.get('experience', 'A renseigner')}\n\n"
            f"Diplomes :\n{self.promoter.get('diplomas', 'A renseigner')}\n\n"
            f"Motivations :\n{self.promoter.get('motivations', 'Creation de valeur et developpement regional')}"
        )
        self._body(slide, info, Inches(3.3), Inches(1.2), Inches(9), Inches(4.5), size=16)
        self._footer(slide)

    def slide04_key_facts(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Le projet en un coup d'oeil")
        inv = self.ctx.results.totalInvestment
        jobs = sum(p.headcount for p in self.ctx.inputs.plAssumptions.personnel)
        ca1 = self.ctx.rev()[0] if self.ctx.rev() else 0
        start = str(self.plan_data.get("start_date") or datetime.now().strftime("%m/%Y"))
        cards = [
            ("💼", "Secteur", self.sector),
            ("📍", "Localisation", self.site),
            ("💰", "Investissement", f"{_fmt(inv)} DT"),
            ("👥", "Emplois", str(jobs or "—")),
            ("📅", "Demarrage", start),
            ("📊", "CA an 1", f"{_fmt(ca1)} DT"),
        ]
        positions = [
            (Inches(0.5), Inches(1.3)),
            (Inches(4.5), Inches(1.3)),
            (Inches(8.5), Inches(1.3)),
            (Inches(0.5), Inches(3.5)),
            (Inches(4.5), Inches(3.5)),
            (Inches(8.5), Inches(3.5)),
        ]
        for (l, t, v), (x, y) in zip(cards, positions):
            self._kpi_card(slide, x, y, Inches(3.6), Inches(1.8), l, t, v)
        self._table_caption(
            slide,
            "Synthese du projet — investissement, emplois et CA an 1",
            top=Inches(5.6),
        )
        self._footer(slide)

    def slide05_market(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Un marche en croissance")
        cols = [
            ("Taille du marche", str(self.market.get("market_size", "Marche national en expansion"))),
            ("Tendances", str(self.market.get("trends", "Demande croissante produits valorises"))),
            ("Part visee", str(self.market.get("share", "Segment qualitatif / regional"))),
        ]
        x = Inches(0.5)
        for title, txt in cols:
            self._body(slide, title, x, Inches(1.2), Inches(3.8), Inches(0.35), size=16, bold=True)
            self._body(slide, str(txt)[:200], x, Inches(1.55), Inches(3.8), Inches(1.2), size=14)
            x += Inches(4.2)
        ipc = self.market.get("ipc_table") or {}
        if ipc:
            p = self.chart_dir / "ipc.png"
            _chart_ipc_mini(ipc if isinstance(ipc, dict) else {}, p)
            self._add_picture(slide, str(p), Inches(0.5), Inches(3.2), Inches(4))
            self._figure_caption(slide, "Evolution IPC du secteur", Inches(0.5), Inches(5.9), width=Inches(4.5))
        quote = self.market.get("quote") or "La consommation evolue vers des produits prêts a l'emploi de qualite."
        self._body(slide, f'"{quote}"', Inches(5), Inches(4.5), Inches(7.5), Inches(0.8), size=14, italic=True)
        self._footer(slide)

    def slide06_swot(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Analyse SWOT")
        defaults = {
            "forces": ["Innovation produit", "Variety", "Rentabilite", "Large clientele"],
            "weaknesses": ["Concurrence importee"],
            "opportunities": ["Peu de concurrence locale", "Sensibilite qualite"],
            "threats": ["Entree nouveaux competiteurs"],
        }
        data = {k: self.swot.get(k) or defaults[k] for k in defaults}
        w, h = Inches(5.8), Inches(2.2)
        gap = Inches(0.15)
        ox, oy = Inches(0.55), Inches(1.15)
        quads = [
            ("FORCES", data["forces"], RGBColor(0xC8, 0xE6, 0xC9), ox, oy),
            ("FAIBLESSES", data["weaknesses"], RGBColor(0xFF, 0xE0, 0xB2), ox + w + gap, oy),
            ("OPPORTUNITES", data["opportunities"], RGBColor(0xBB, 0xDE, 0xFB), ox, oy + h + gap),
            ("MENACES", data["threats"], RGBColor(0xFF, 0xCD, 0xD2), ox + w + gap, oy + h + gap),
        ]
        for title, items, fill, x, y in quads:
            sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
            bullets = "\n".join(f"• {it}" for it in items[:4])
            self._body(slide, f"{title}\n{bullets}", x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h - Inches(0.2), size=13)
        self._footer(slide)

    def slide07_products(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Produits & offre commerciale")
        products = _product_names(self.ctx)
        rev = self.ctx.rev()
        n = max(1, len(products))
        x = Inches(0.4)
        for pname in products[:4]:
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(3), Inches(2.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = C_GREY
            part = _pct(_safe_div(rev[0] / n, rev[0])) if rev[0] else "—"
            price = self.ctx.inputs.operations.salePrice
            self._body(
                slide,
                f"{pname}\n\nPrix : {_fmt(price)} DT\nPart CA est. : {part}\n\n[Image]",
                x + Inches(0.1),
                Inches(1.35),
                Inches(2.8),
                Inches(2.5),
                size=13,
            )
            x += Inches(3.15)
        self._footer(slide)

    def slide08_investment(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Plan d'investissement")
        buckets = _categorize_investments(self.ctx)
        items: list[tuple[str, float]] = []
        for key, label in [
            ("incorporel", "Incorporelles"),
            ("agencement", "Agencements"),
            ("industriel", "Corporelles / industriel"),
            ("transport", "Transport"),
            ("bureau", "Bureau"),
            ("preliminaires", "Frais preliminaires"),
        ]:
            for name, amt, *_ in buckets.get(key, []):
                items.append((f"{label}: {name[:20]}", amt))
        bfr = self.ctx.results.bfr.years[0] if self.ctx.results.bfr.years else 0
        items.append(("BFR initial", bfr))
        total = self.ctx.results.totalInvestment + bfr
        p = self.chart_dir / "invest.png"
        _chart_investment_bars(items, total, p)
        self._add_picture(slide, str(p), Inches(0.4), Inches(1.1), Inches(8.5))
        self._figure_caption(slide, "Repartition du plan d'investissement (DT)", Inches(0.4), Inches(5.85))
        self._body(
            slide,
            f"TOTAL : {_fmt(total)} DT",
            Inches(9), Inches(2.5), Inches(3.8), Inches(1),
            size=32,
            bold=True,
            color=C_PRIMARY,
        )
        self._footer(slide)

    def slide09_financing(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Schema de financement")
        fin = self.ctx.inputs.financing
        total = self.ctx.results.totalInvestment
        eq = total * fin.equityRatio
        debt = fin.loan.amount or total * fin.debtRatio
        sub = float(self.plan_data.get("subsidies") or 0)
        p = self.chart_dir / "financing.png"
        _chart_financing_pie(
            ["Fonds propres", "CMT", "Subventions"],
            [eq, debt, sub if sub else 0.01],
            p,
        )
        self._add_picture(slide, str(p), Inches(0.4), Inches(1.1), Inches(4.5))
        self._figure_caption(slide, "Schema de financement (fonds propres, dette, subventions)", Inches(0.4), Inches(5.5), width=Inches(4.8))
        loan = fin.loan
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(7.5), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = C_GREY
        txt = (
            f"Montant : {_fmt(debt)} DT\n"
            f"Taux : {_pct(loan.rate)}\n"
            f"Duree : {loan.years} ans\n"
            f"Differe principal : {loan.graceMonthsPrincipal} mois"
        )
        self._body(slide, "Conditions du credit\n\n" + txt, Inches(5.4), Inches(1.35), Inches(7), Inches(2), size=15)
        rows = [
            ["Fonds propres", _fmt(eq), _pct(fin.equityRatio)],
            ["Credit MLT", _fmt(debt), _pct(fin.debtRatio)],
            ["Subventions", _fmt(sub), "—"],
        ]
        y = Inches(3.6)
        for row in rows:
            self._body(slide, " | ".join(row), Inches(5.2), y, Inches(7.5), Inches(0.35), size=14)
            y += Inches(0.4)
        self._footer(slide)

    def slide10_revenue(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Chiffre d'affaires previsionnel")
        rev = self.ctx.rev()
        growth = [0.0] + [_safe_div(rev[i] - rev[i - 1], rev[i - 1]) for i in range(1, HORIZON)]
        p = self.chart_dir / "ca.png"
        _chart_ca_bars(rev, growth, p)
        self._add_picture(slide, str(p), Inches(0.5), Inches(1.05), Inches(12))
        self._figure_caption(slide, "Chiffre d'affaires previsionnel sur 7 ans (DT)", Inches(0.5), Inches(5.75))
        cap = self.plan_data.get("capacity_pct") or "93"
        self._body(
            slide,
            f"Pleine capacite a partir de l'annee de stabilisation : {cap}%",
            Inches(0.5),
            Inches(6.2),
            Inches(11),
            Inches(0.4),
            size=14,
            italic=True,
        )
        self._footer(slide)

    def slide11_results(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Evolution des resultats")
        ebit = []
        for y in range(HORIZON):
            yr = self.ctx.yearly[y] if y < len(self.ctx.yearly) else {}
            ebit.append(float(yr.get("ebe", 0)) - float(yr.get("depreciation", 0)))
        p = self.chart_dir / "results.png"
        _chart_results_lines(self.ctx.rev(), ebit, self.ctx.net(), p)
        self._add_picture(slide, str(p), Inches(0.4), Inches(1.05), Inches(12))
        self._figure_caption(slide, "Evolution CA, resultat d'exploitation et resultat net", Inches(0.4), Inches(5.75))
        self._footer(slide)

    def slide12_kpis(self) -> None:
        slide = self._blank_slide()
        title = "Rentabilite — indicateurs cles"
        if self.audience == "banque":
            title = "Rentabilite & capacite de remboursement"
        elif self.audience == "investisseur":
            title = "Rentabilite & creation de valeur"
        self._title_box(slide, title)
        ind = self.ctx.ind
        loan_rate = self.ctx.inputs.financing.loan.rate
        margin = 0.0
        if self.ctx.yearly:
            margin = sum(float(y.get("grossMarginPct", 0)) for y in self.ctx.yearly) / len(self.ctx.yearly)
        kpis = [
            ("TRI", _pct(ind.tri), f"vs taux banque {_pct(loan_rate)}"),
            ("VAN", f"{_fmt(ind.van)} DT", "Creation de valeur positive" if ind.van >= 0 else "A ameliorer"),
            ("DRCI", _drci_label(ind.drciYears), "Recuperation investissement"),
            ("Marge brute", _pct(margin), "Moyenne sur 7 ans"),
        ]
        if self.audience == "investisseur":
            kpis[0] = ("TRI", _pct(ind.tri), "Rendement investisseur cible")
            kpis[1] = ("VAN", f"{_fmt(ind.van)} DT", "Potentiel de croissance")
        elif self.audience == "client":
            kpis = [
                ("Emplois", str(sum(p.headcount for p in self.ctx.inputs.plAssumptions.personnel)), "Impact local"),
                ("CA an 7", f"{_fmt(self.ctx.rev()[-1])} DT", "Developpement commercial"),
                ("Produits", str(len(_product_names(self.ctx))), "Offre diversifiee"),
                ("Marge", _pct(margin), "Qualite economique"),
            ]
        x = Inches(0.45)
        for label, val, sub in kpis:
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.15), Inches(3), Inches(2.5)
            )
            card.line.color.rgb = C_ACCENT if label == "VAN" else C_PRIMARY
            self._fill_kpi_card(card, label, val, sub)
            x += Inches(3.15)
        self._table_caption(
            slide,
            "Indicateurs de rentabilite (TRI, VAN, DRCI, marge) — lecture banque / investisseur",
            top=Inches(4.0),
        )
        self._footer(slide)

    def slide13_cashflow(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Cash-flows & remboursement")
        loan_rows = _loan_annual_rows(self.ctx)
        p = self.chart_dir / "cashflow.png"
        _chart_cashflow_debt(self.ctx.ocf(), loan_rows, p)
        self._add_picture(slide, str(p), Inches(0.4), Inches(1.05), Inches(9))
        self._figure_caption(slide, "Cash-flows d'exploitation et capital restant du", Inches(0.4), Inches(5.7), width=Inches(9))
        msg = "Capacite de remboursement demontree"
        if self.audience == "investisseur":
            msg = "Tresorerie soutenant la croissance"
        self._body(slide, msg, Inches(9.8), Inches(2.5), Inches(3), Inches(1), size=20, bold=True, color=C_PRIMARY)
        self._footer(slide)

    def slide14_loan_table(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Tableau d'amortissement (resume)")
        loan_rows = _loan_annual_rows(self.ctx)
        headers = ["An", "Capital restant", "Interets", "Principal", "Service dette"]
        y = Inches(1.15)
        col_w = [Inches(0.6), Inches(2.2), Inches(2), Inches(2), Inches(2.2)]
        x_cols = [Inches(0.4)]
        for w in col_w[:-1]:
            x_cols.append(x_cols[-1] + w)
        for i, h in enumerate(headers):
            self._body(slide, h, x_cols[i], y, col_w[i], Inches(0.35), size=12, bold=True)
        y += Inches(0.4)
        total_int = 0.0
        grace = self.ctx.inputs.financing.loan.graceMonthsPrincipal
        for yi, row in enumerate(loan_rows):
            total_int += row["interest"]
            line = f"An {yi + 1}" + (" *" if yi == 0 and grace else "")
            vals = [
                line,
                _fmt(row["closing"]),
                _fmt(row["interest"]),
                _fmt(row["principal"]),
                _fmt(row["service"]),
            ]
            for i, v in enumerate(vals):
                self._body(slide, v, x_cols[i], y, col_w[i], Inches(0.32), size=11)
            y += Inches(0.32)
        self._body(slide, f"Total interets : {_fmt(total_int)} DT", Inches(0.4), y + Inches(0.15), Inches(6), Inches(0.4), size=14, bold=True)
        if grace:
            self._body(slide, "* Periode de grace sur principal", Inches(0.4), y + Inches(0.45), Inches(5), Inches(0.3), size=11, italic=True)
        self._table_caption(slide, "Tableau d'amortissement annuel (resume)", top=y + Inches(0.75))
        self._footer(slide)

    def slide15_costs(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Structure des couts")
        personnel = sum(p.headcount * p.annualSalary for p in self.ctx.inputs.plAssumptions.personnel) * HORIZON
        achats = sum(self.ctx.purchase_mp())
        dap = sum(self.ctx.dep())
        autres = sum(self.ctx.mkt()) + sum(self.ctx.dist())
        vals = [achats, personnel, dap, autres]
        labels = ["Achats matieres", "Personnel", "Amortissements", "Autres"]
        var_pct = _pct(_safe_div(achats, sum(vals)))
        p = self.chart_dir / "costs.png"
        _chart_cost_donut(vals, labels, p)
        self._add_picture(slide, str(p), Inches(0.5), Inches(1.05), Inches(6))
        self._figure_caption(slide, "Structure des couts (moyenne 7 ans)", Inches(0.5), Inches(5.7), width=Inches(6.5))
        self._body(
            slide,
            f"Structure maitrisee — {var_pct} de couts variables (achats)",
            Inches(7),
            Inches(2.5),
            Inches(5.5),
            Inches(0.8),
            size=16,
        )
        self._footer(slide)

    def slide16_breakeven(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Point mort & sensibilite")
        rev = self.ctx.rev()
        fixed = 0.0
        var = 0.0
        for y in range(HORIZON):
            yr = self.ctx.yearly[y] if y < len(self.ctx.yearly) else {}
            cogs = float(yr.get("cogs", 0))
            var += cogs
            fixed += float(yr.get("personnel", 0)) + float(yr.get("depreciation", 0)) + float(yr.get("otherOpex", 0))
        var_avg = var / HORIZON
        fix_avg = fixed / HORIZON
        rev_avg = sum(rev) / HORIZON
        margin_rate = _safe_div(rev_avg - var_avg, rev_avg)
        be = _safe_div(fix_avg, margin_rate) if margin_rate else 0
        _mpl_readability()
        fig, ax = plt.subplots(figsize=(6, 4))
        x = [0, rev_avg * 1.5]
        ax.plot(x, [fix_avg + margin_rate * xi for xi in x], color=MPL_PRIMARY, label="Couts totaux")
        ax.plot(x, x, color=MPL_GREEN, label="CA")
        ax.axvline(be, color=MPL_ORANGE, linestyle="--", label="Point mort")
        ax.legend(fontsize=8)
        ax.set_title("Seuil de rentabilite", fontsize=10, color=MPL_PRIMARY)
        p = self.chart_dir / "breakeven.png"
        _save_chart(fig, p)
        self._add_picture(slide, str(p), Inches(0.4), Inches(1.1), Inches(6))
        self._figure_caption(slide, "Seuil de rentabilite (point mort)", Inches(0.4), Inches(5.5), width=Inches(6.5))
        sens = _sensitivity_mini(self.ctx)
        self._body(slide, "VAN (±10% CA / couts)", Inches(7), Inches(1.2), Inches(5.5), Inches(0.35), size=14, bold=True)
        hdr = "     ".join(["CA 90%", "100%", "110%"])
        self._body(slide, hdr, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3), size=11)
        for i, cost_l in enumerate(["110%", "100%", "90%"]):
            self._body(slide, f"C {cost_l}  " + "  ".join(sens[i]), Inches(7), Inches(1.95 + i * 0.32), Inches(5.5), Inches(0.3), size=10)
        self._body(
            slide,
            f"Point mort : {_fmt(be)} DT ({_pct(_safe_div(be, rev_avg))} capacite)",
            Inches(7),
            Inches(3.2),
            Inches(5.5),
            Inches(0.5),
            size=13,
        )
        self._footer(slide)

    def slide17_planning(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Planning de realisation")
        phases = self.plan_data.get("planning") or [
            {"phase": "Etude & dossier", "months": 3, "color": MPL_PRIMARY},
            {"phase": "Financement", "months": 5, "color": MPL_SECONDARY},
            {"phase": "Commande equipements", "months": 5, "color": MPL_TEAL},
            {"phase": "Demarrage production", "months": 4, "color": MPL_GREEN},
        ]
        ox, oy, h = Inches(0.6), Inches(1.5), Inches(0.55)
        scale = Inches(0.35)
        x = ox
        if isinstance(phases, list) and phases and isinstance(phases[0], dict):
            for ph in phases:
                w = scale * int(ph.get("months", 3))
                sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, oy, w, h)
                sh.fill.solid()
                sh.fill.fore_color.rgb = C_SECONDARY
                self._body(slide, f"{ph.get('phase', '')} ({ph.get('months', 3)} mois)", x, oy + h + Inches(0.08), w, Inches(0.5), size=11)
                x += w + Inches(0.08)
        self._footer(slide)

    def slide18_legal(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Aspects juridiques & avantages fiscaux")
        law = str(
            self.plan_data.get("legal_framework")
            or "Code des incitations aux investissements — Tunisie"
        )
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.15), Inches(12), Inches(0.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = C_GREY
        self._body(slide, law, Inches(0.65), Inches(1.25), Inches(11.5), Inches(0.6), size=14)
        rows = [
            ["Prime investissement", str(self.plan_data.get("prime_invest") or "10% equipements (plafond)")],
            ["Prime etude", str(self.plan_data.get("prime_etude") or "70% cout etude (plafond)")],
            ["Prime immateriels", str(self.plan_data.get("prime_immat") or "50% immateriels")],
        ]
        y = Inches(2.2)
        for row in rows:
            self._body(slide, f"{row[0]} : {row[1]}", Inches(0.6), y, Inches(11.5), Inches(0.45), size=15)
            y += Inches(0.5)
        self._footer(slide)

    def slide19_team(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Notre engagement & equipe")
        self._body(
            slide,
            "Une equipe engagee pour la reussite industrielle et commerciale du projet.",
            Inches(0.5),
            Inches(1.1),
            Inches(12),
            Inches(0.5),
            size=16,
            italic=True,
        )
        team = self.team if isinstance(self.team, list) and self.team else []
        if not team:
            for p in self.ctx.inputs.plAssumptions.personnel[:6]:
                if p.role.strip():
                    team.append({"role": p.role, "name": "A definir", "profile": f"{p.headcount} poste(s)"})
        y = Inches(1.8)
        for member in team[:5]:
            if isinstance(member, dict):
                line = f"• {member.get('role', 'Poste')} — {member.get('name', '')} : {member.get('profile', '')}"
            else:
                line = f"• {member}"
            self._body(slide, line, Inches(0.6), y, Inches(11.5), Inches(0.45), size=15)
            y += Inches(0.5)
        self._footer(slide)

    def slide20_conclusion(self) -> None:
        slide = self._blank_slide()
        self._title_box(slide, "Conclusion")
        ind = self.ctx.ind
        if self.audience == "banque":
            bullets = [
                f"Capacite de remboursement : DRCI {_drci_label(ind.drciYears)}",
                f"Service de dette couvert par les cash-flows",
                f"Projet crediteur : VAN {_fmt(ind.van)} DT",
            ]
        elif self.audience == "investisseur":
            bullets = [
                f"TRI attractif : {_pct(ind.tri)}",
                f"VAN : {_fmt(ind.van)} DT — potentiel de valeur",
                "Marche en croissance et offre differenciante",
            ]
        else:
            bullets = [
                f"{sum(p.headcount for p in self.ctx.inputs.plAssumptions.personnel)} emplois et ancrage local",
                f"CA an 1 : {_fmt(self.ctx.rev()[0])} DT",
                "Produits de qualite pour le marche tunisien",
            ]
        y = Inches(1.3)
        for b in bullets:
            self._body(slide, b, Inches(0.7), y, Inches(11.5), Inches(0.7), size=22, bold=True)
            y += Inches(0.85)
        contact = self.plan_data.get("contact") or self.promoter.get("contact") or "contact@projet.tn"
        self._body(slide, f"Contact : {contact}", Inches(0.7), Inches(5.5), Inches(8), Inches(0.5), size=16)
        self._body(slide, "Merci pour votre confiance", Inches(0.7), Inches(6.1), Inches(8), Inches(0.5), size=20, bold=True, color=C_PRIMARY)
        qr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.5), Inches(5.2), Inches(1.5), Inches(1.5))
        qr.fill.solid()
        qr.fill.fore_color.rgb = C_GREY
        self._body(slide, "[QR]", Inches(10.85), Inches(5.75), Inches(1), Inches(0.4), size=12)
        self._footer(slide)

    def build(self, output_path: str) -> None:
        self.chart_dir = Path(tempfile.mkdtemp(prefix="bp_pptx_"))
        try:
            self.slide01_cover()
            self.slide02_agenda()
            self.slide03_promoter()
            self.slide04_key_facts()
            self.slide05_market()
            self.slide06_swot()
            self.slide07_products()
            self.slide08_investment()
            self.slide09_financing()
            self.slide10_revenue()
            self.slide11_results()
            self.slide12_kpis()
            self.slide13_cashflow()
            self.slide14_loan_table()
            self.slide15_costs()
            self.slide16_breakeven()
            self.slide17_planning()
            self.slide18_legal()
            self.slide19_team()
            self.slide20_conclusion()
            out = Path(output_path)
            out.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(str(out))
        finally:
            if self.chart_dir:
                shutil.rmtree(self.chart_dir, ignore_errors=True)


def generate_pptx_presentation(
    plan_data: dict,
    audience: str,
    output_path: str,
) -> None:
    """
    Generate a 20-slide PowerPoint deck from plan_data.

    audience: 'banque' | 'investisseur' | 'client'
    """
    if audience not in ("banque", "investisseur", "client"):
        raise ValueError("audience must be 'banque', 'investisseur' or 'client'")
    deck = PptxDeck(plan_data, audience)  # type: ignore[arg-type]
    deck.build(output_path)


def celery_export_pptx(
    plan_data: dict,
    audience: str,
    output_path: str,
) -> str:
    generate_pptx_presentation(plan_data, audience, output_path)
    return str(Path(output_path).resolve())


def build_pptx_from_plan(
    plan_data: dict,
    export_dir: Path,
    *,
    audience: str = "banque",
) -> str:
    ctx = _parse_plan_data(plan_data)
    slug = "".join(c if c.isalnum() else "_" for c in ctx.company)[:40].strip("_") or "projet"
    path = export_dir / f"presentation_{ctx.plan_id}_{slug}.pptx"
    generate_pptx_presentation(plan_data, audience, str(path))
    return str(path.resolve())
